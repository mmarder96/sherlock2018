import PyPDF2
from pptx import Presentation
from io import StringIO
import docx


def read_pdf(filename):
    # read doc
    pdf = PyPDF2.PdfFileReader(filename)
    text = []
    numPages = pdf.getNumPages()
    for page in range(0, numPages):
        pageObj = pdf.getPage(page)
        content = pageObj.extractText().encode('utf-8', 'ignore')
        text.append(content)

    # combine/split elements in list
    text_by_limit = []
    upload_length = 0
    max_length = 5000
    buf = StringIO()
    for page in text:
        page_length = len(page)
        if upload_length + page_length < max_length:
            raw_text = str(page).split('\\n')
            for raw in raw_text:
                if raw == '' or raw == ' ':
                    continue
                line = str(raw).rstrip()

                if not line.endswith(':') \
                        and not line.endswith(',') \
                        and not line.endswith('.') \
                        and not line.endswith('…'):
                    line += '. '

                buf.write(line)
            upload_length += page_length
        else:
            text_by_limit.append(buf.getvalue())
            buf = StringIO()
            upload_length = 0
    buf.close()

    return text_by_limit


def read_pptx(filename):
    prs = Presentation(filename)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                raw_text = shape.text.split('\n')

                for raw in raw_text:
                    if raw == '':
                        continue
                    line = str(raw).rstrip()

                    if not line.endswith(':') \
                            and not line.endswith(',') \
                            and not line.endswith('.') \
                            and not line.endswith('…'):
                        line += '.'

                    text.append(line)
            except AttributeError:
                # Placeholder object has no text
                pass

    # strip text of empty spaces
    stripped_text = []
    for line in text:
        stripped_text.append(str(line).strip())

    return '\n'.join(stripped_text)


def read_docx(filename):
    # read doc
    doc = docx.Document(filename)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)

    # combine/split elements in list,
    # so that each element is less than the limit
    text_by_limit = []
    upload_size = 0
    max_byte_size = 5000
    buf = StringIO()
    for para in text:
        para_size = len(para.encode('utf-8'))
        if upload_size + para_size < max_byte_size:
            raw_text = str(para).split('\n')
            for raw in raw_text:
                if raw == '':
                    continue

                line = str(raw).rstrip()
                if not line.endswith(':') \
                        and not line.endswith(',') \
                        and not line.endswith('.') \
                        and not line.endswith('…'):
                    line += '. '

                buf.write(line)
            upload_size += para_size
        else:
            text_by_limit.append(buf.getvalue())
            buf = StringIO()
            upload_size = 0
    buf.close()

    return text_by_limit


def write_pptx(pptx_output_str):
    with open('pptx-test', 'w') as f:
        f.write(pptx_output_str)


def write_docx(docx_output_list):
    with open('docx-test', 'w') as f:
        for text in docx_output_list:
            f.write(str(text.encode('utf-8')) + '\n')


if __name__ == '__main__':
    # TODO replace this
    pdf_filename = 'Lecture11.pdf'
    pptx_filename = 'graphics.pptx'
    docx_filename = 'homework9.docx'

    read_pdf(pdf_filename)
    pptx_output_str = read_pptx(pptx_filename)
    docx_output_list = read_docx(docx_filename)

    write_pptx(pptx_output_str)
    write_docx(docx_output_list)
